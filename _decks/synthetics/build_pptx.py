"""Build Synthetics Deck.pptx — native PowerPoint version of the HTML deck.

Matches the deck's 1920x1080 design grid by mapping pixels -> EMU at 6350 EMU/px
so the slide is 12192000 x 6858000 EMU (PowerPoint's standard 16:9 widescreen).
"""
from pathlib import Path
from pptx import Presentation
from pptx.util import Emu, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from lxml import etree

HERE = Path(__file__).parent
ROOT = HERE.parent.parent
OUT = HERE / "Synthetics Deck.pptx"
PORTRAIT = ROOT / "assets" / "img" / "me_sketch_cleaned.png"

# ── design grid → EMU ────────────────────────────────────────────
PX = 6350  # EMU per design pixel (1920px → 12192000 EMU = 13.333in)
def E(px): return Emu(int(px * PX))

# ── palette (from CSS) ───────────────────────────────────────────
BG_PRIMARY   = RGBColor(0x0a, 0x0e, 0x1a)
BG_SECONDARY = RGBColor(0x11, 0x18, 0x27)
BG_CARD_RGB  = RGBColor(0x14, 0x18, 0x24)   # ≈ rgba(255,255,255,0.04) on bg
BG_GLASS_RGB = RGBColor(0x1a, 0x1e, 0x2a)
TEXT_PRIMARY = RGBColor(0xe8, 0xe6, 0xe3)
TEXT_SECOND  = RGBColor(0x9c, 0xa3, 0xaf)
TEXT_MUTED   = RGBColor(0x6b, 0x72, 0x80)
ACCENT       = RGBColor(0xc9, 0xa9, 0x6e)
ACCENT_LIGHT = RGBColor(0xdf, 0xc5, 0x99)
ROSE         = RGBColor(0xf4, 0x3f, 0x5e)
GREEN        = RGBColor(0x10, 0xb9, 0x81)
BLUE         = RGBColor(0x7e, 0xb8, 0xda)
BORDER_RGB   = RGBColor(0x2a, 0x2e, 0x3a)

FONT_BODY    = "Inter"
FONT_DISPLAY = "Outfit"
FONT_MONO    = "JetBrains Mono"

NOTES = [
    "A heat engine was in serious industrial use for about a century before anyone had a theory of what it was doing. Before Carnot and Clausius, the engine was a useful mystery — you could build one, and if you were careful and lucky, it worked. After the theory, refrigeration, systematic engine design, and whole industries nobody had been able to imagine became routine. The theory did not invent the engine; it made the engine legible, and then it made everything downstream of the engine possible.",
    "My claim is that we are today in an exactly analogous situation with intelligence. Intelligence engines exist and work — biological brains, deep networks, the machinery of modern AI. But a principled physical theory of what they are, what they cost, and what any given substrate can efficiently do, does not. We are the mechanics before the physicists. Synthetics is my attempt at that theory.",
    "Two foundational pre-preprints already lay the static mathematics. The Cost of Complexity derives, from three sparse assumptions, the thermodynamic price of any embodied representation — and identifies a sharp critical ceiling where the partition function becomes the Riemann zeta function. Equilibrium from the Inside Out gives an exact algebraic closure theorem — a criterion for when a subsystem admits a compact local effective description — proven for a genuinely non-commuting case. Both are currently static. The residency extends them into dynamics, connects them formally, and applies them to neuromorphic computing.",
    "Here is the unifying move. Biological cortex, open quantum systems, and neuromorphic hardware look nothing like each other as physical objects, but they are all dynamical systems, and dynamical systems theory lets us write them all down as the same linear evolution over a lifted observable space — one operator algebra for all three. That shared representation is the glue. And the moment you insist these dynamics be physically embodied, two constraints fall out whether you want them to or not — a thermodynamic cost, and an algebraic closure criterion. Design principles for efficient intelligence are just what those constraints look like when you take them seriously. That is the shape of Synthetics.",
    "Alongside the mathematics, public writing is treated as part of the work rather than a separate communications activity. The argument is simple: a synthesis which cannot be stated across the fields it unifies has not unified them. Explicability is a research deliverable, not a chore. Little Learning Machine and the FQxI essay are where that work is already happening — second place out of several hundred entrants for an essay on whether life comes from outer Hilbert space.",
    "A word on who's doing this. I'm a theoretical physicist — PhD King's, postdoc at Tulane — sitting at the intersection of the three fields Synthetics wants to unify: quantum theory, thermodynamics, and machine intelligence. Twenty-five-plus papers in quantum control, open systems, reservoir computing, superoscillations, Koopman methods. The website holds the full catalogue — the pre-preprints at slash physics, the essays on Substack, and the broader writing at slash writing. Thank you.",
]

# ── helpers ──────────────────────────────────────────────────────
def set_slide_bg(slide, color):
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = color

def add_rect(slide, x, y, w, h, fill=None, line=None, line_w=0.75, radius=False):
    shape_type = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    sh = slide.shapes.add_shape(shape_type, E(x), E(y), E(w), E(h))
    if radius:
        # set adjustment for corner radius — fraction of shorter side
        sh.adjustments[0] = 0.06
    if fill is None:
        sh.fill.background()
    else:
        sh.fill.solid()
        sh.fill.fore_color.rgb = fill
    if line is None:
        sh.line.fill.background()
    else:
        sh.line.color.rgb = line
        sh.line.width = Pt(line_w)
    sh.shadow.inherit = False
    return sh

def add_text(slide, x, y, w, h, text, *, font=FONT_BODY, size=18, color=TEXT_PRIMARY,
             bold=False, italic=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP,
             tracking=None):
    tb = slide.shapes.add_textbox(E(x), E(y), E(w), E(h))
    tf = tb.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    # text may be a string OR list of (text, opts) runs
    if isinstance(text, str):
        runs = [(text, {})]
    else:
        runs = text
    p = tf.paragraphs[0]
    p.alignment = align
    for i, (t, opts) in enumerate(runs):
        if t == "\n":
            p = tf.add_paragraph()
            p.alignment = align
            continue
        r = p.add_run()
        r.text = t
        rf = r.font
        rf.name = opts.get("font", font)
        rf.size = Pt(opts.get("size", size))
        rf.bold = opts.get("bold", bold)
        rf.italic = opts.get("italic", italic)
        rf.color.rgb = opts.get("color", color)
        if tracking:
            # letter-spacing in 100ths of a point — best-effort
            rPr = r._r.get_or_add_rPr()
            rPr.set("spc", str(tracking))
    return tb

def add_kicker(slide, x, y, label):
    # gold rule + uppercase tracked label
    add_rect(slide, x, y + 14, 48, 1, fill=ACCENT)
    add_text(slide, x + 62, y, 1200, 36, label.upper(),
             font=FONT_BODY, size=18, color=ACCENT, bold=True, tracking=200)

def add_footer(slide, brand="Synthetics", meta="Gerard McCaul — Astera residency application"):
    add_text(slide, 120, 1014, 1200, 36,
             [(brand, {"font": FONT_DISPLAY, "size": 20, "bold": True, "color": TEXT_MUTED}),
              (".", {"font": FONT_DISPLAY, "size": 20, "bold": True, "color": ACCENT}),
              ("   " + meta, {"font": FONT_BODY, "size": 16, "color": TEXT_MUTED})])

def add_counter(slide, n, total=6):
    add_text(slide, 1620, 1014, 200, 36, f"{n:02d} / {total:02d}",
             font=FONT_MONO, size=16, color=TEXT_MUTED, align=PP_ALIGN.RIGHT, tracking=100)

def add_speaker_note(slide, text):
    notes_tf = slide.notes_slide.notes_text_frame
    notes_tf.text = text

# ── build ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width = E(1920)
prs.slide_height = E(1080)
blank = prs.slide_layouts[6]

# ───── SLIDE 1 — Carnot title ─────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG_PRIMARY)
# era divider top-right
add_text(s, 1320, 44, 480, 32, "C. 1712 — 1824",
         font=FONT_MONO, size=16, color=TEXT_MUTED, align=PP_ALIGN.RIGHT, tracking=120)
add_text(s, 1320, 78, 480, 32, "ENGINE BEFORE THEORY",
         font=FONT_MONO, size=16, color=ACCENT, align=PP_ALIGN.RIGHT, tracking=180, bold=True)

add_kicker(s, 160, 270, "Synthetics — Residency Pitch")
# Title (3 lines)
add_text(s, 160, 320, 1600, 360, [
    ("The engine ran\n", {"font": FONT_DISPLAY, "size": 90, "bold": True, "color": TEXT_PRIMARY}),
    ("for a ", {"font": FONT_DISPLAY, "size": 90, "bold": True, "color": TEXT_PRIMARY}),
    ("century", {"font": FONT_DISPLAY, "size": 90, "bold": True, "color": ACCENT}),
    ("\n", {}),
    ("before the theory.", {"font": FONT_DISPLAY, "size": 90, "bold": True, "color": TEXT_PRIMARY}),
], align=PP_ALIGN.LEFT)

# Stat row
add_text(s, 160, 750, 760, 140, "~112",
         font=FONT_DISPLAY, size=140, bold=True, color=ACCENT)
add_text(s, 160, 900, 760, 80,
         "years between Newcomen's atmospheric engine and Carnot's Réflexions.",
         font=FONT_BODY, size=18, color=TEXT_SECOND)

add_text(s, 980, 760, 800, 130,
         [("useful mystery\n", {"font": FONT_DISPLAY, "size": 40, "bold": True, "color": TEXT_PRIMARY}),
          ("→ routine industry", {"font": FONT_DISPLAY, "size": 40, "bold": True, "color": TEXT_PRIMARY})])
add_text(s, 980, 920, 800, 60,
         "Refrigeration, systematic engine design, entire industries nobody had been able to imagine.",
         font=FONT_BODY, size=18, color=TEXT_SECOND)

add_footer(s); add_counter(s, 1)
add_speaker_note(s, NOTES[0])

# ───── SLIDE 2 — We are here again ─────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG_PRIMARY)
add_kicker(s, 120, 130, "The claim")
add_text(s, 120, 180, 1600, 110, "We are here again.",
         font=FONT_DISPLAY, size=72, bold=True, color=TEXT_PRIMARY)
add_text(s, 120, 310, 1500, 130,
         [("Intelligence engines exist and work.\n", {"font": FONT_BODY, "size": 28, "color": TEXT_SECOND}),
          ("A physical theory of them does not.", {"font": FONT_BODY, "size": 28, "color": TEXT_SECOND})])

# Two cards
def twin_col(x, exists=True, label=None, h2=None, p=None, items=None):
    border = ACCENT if exists else BORDER_RGB
    add_rect(s, x, 510, 800, 440, fill=BG_CARD_RGB, line=border, line_w=1.2, radius=True)
    add_text(s, x + 36, 540, 720, 28,
             ("■ " if exists else "□ ") + label,
             font=FONT_MONO, size=16, color=ACCENT if exists else ROSE, bold=True, tracking=180)
    add_text(s, x + 36, 580, 720, 70, h2,
             font=FONT_DISPLAY, size=42, bold=True, color=TEXT_PRIMARY)
    add_text(s, x + 36, 660, 720, 130, p, font=FONT_BODY, size=18, color=TEXT_SECOND)
    for i, it in enumerate(items):
        add_text(s, x + 36, 800 + i*36, 28, 28, "—",
                 font=FONT_BODY, size=18, color=TEXT_MUTED)
        add_text(s, x + 70, 800 + i*36, 700, 28, it,
                 font=FONT_BODY, size=18, color=TEXT_SECOND)

twin_col(120, exists=True, label="WHAT EXISTS",
         h2="The engines.",
         p="Biological brains. Deep networks. Neuromorphic chips. We build them, operate them, and depend on them — at ruinous energetic cost.",
         items=["They work.", "We have empirical scaling laws.", "We do not know why they work."])
twin_col(1000, exists=False, label="WHAT IS MISSING",
         h2="A theory of them.",
         p="A principled account of what an intelligence is, what it costs, and what any given physical substrate can efficiently do.",
         items=["No Carnot bound.", "No closure criterion.", "No substrate-independent accounting."])

add_footer(s, meta=""); add_counter(s, 2)
add_speaker_note(s, NOTES[1])

# ───── SLIDE 3 — Foundational math ─────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG_PRIMARY)
add_kicker(s, 120, 60, "The foundations are already written")
add_text(s, 120, 110, 1700, 130,
         [("Two pre-preprints,\n", {"font": FONT_DISPLAY, "size": 44, "bold": True, "color": TEXT_PRIMARY}),
          ("one static skeleton.", {"font": FONT_DISPLAY, "size": 44, "bold": True, "color": TEXT_PRIMARY})])
add_text(s, 120, 250, 1700, 70,
         [("The residency extends both into ", {"font": FONT_BODY, "size": 18, "color": TEXT_SECOND}),
          ("dynamics", {"font": FONT_BODY, "size": 18, "color": ACCENT, "bold": True}),
          (", connects them formally, and applies them to ", {"font": FONT_BODY, "size": 18, "color": TEXT_SECOND}),
          ("neuromorphic computing", {"font": FONT_BODY, "size": 18, "color": ACCENT, "bold": True}),
          (".", {"font": FONT_BODY, "size": 18, "color": TEXT_SECOND})])

# Browser window
add_rect(s, 120, 340, 1680, 620, fill=BG_SECONDARY, line=BORDER_RGB, line_w=1, radius=True)
# Window bar
add_rect(s, 120, 340, 1680, 56, fill=BG_GLASS_RGB, line=BORDER_RGB, line_w=1)
for i, c in enumerate([RGBColor(0x40,0x40,0x48)]*3):
    add_rect(s, 142 + i*22, 358, 16, 16, fill=c)
add_text(s, 600, 354, 720, 32,
         [("gmccaul.co.uk", {"font": FONT_MONO, "size": 16, "color": TEXT_SECOND, "bold": True}),
          ("/physics.html  #pre-preprints", {"font": FONT_MONO, "size": 16, "color": TEXT_MUTED})],
         align=PP_ALIGN.CENTER)

# Section label inside
add_rect(s, 158, 437, 24, 1, fill=ACCENT)
add_text(s, 192, 422, 400, 28, "DRAFTS",
         font=FONT_BODY, size=16, bold=True, color=ACCENT, tracking=160)
add_text(s, 158, 450, 1400, 70, "Pre-preprints",
         font=FONT_DISPLAY, size=36, bold=True, color=TEXT_PRIMARY)
add_text(s, 158, 520, 1500, 50,
         "In-progress drafts of a foundational document for a broader program concerned with a physical theory of representations.",
         font=FONT_BODY, size=16, color=TEXT_SECOND)

# Two paper cards
def paper_card(x, title, journal, blurb):
    add_rect(s, x, 600, 800, 320, fill=BG_CARD_RGB, line=ACCENT, line_w=1.2, radius=True)
    add_text(s, x + 28, 620, 740, 60, title,
             font=FONT_DISPLAY, size=24, bold=True, color=TEXT_PRIMARY)
    add_text(s, x + 28, 690, 740, 28,
             [("G. M", {"font": FONT_BODY, "size": 16, "color": ACCENT, "bold": True}),
              ("c", {"font": FONT_BODY, "size": 11, "color": ACCENT, "bold": True}),
              ("Caul", {"font": FONT_BODY, "size": 16, "color": ACCENT, "bold": True})])
    add_text(s, x + 28, 718, 740, 28, journal,
             font=FONT_BODY, size=16, italic=True, color=TEXT_MUTED)
    add_text(s, x + 28, 758, 740, 110, blurb,
             font=FONT_BODY, size=15, color=TEXT_SECOND)
    add_rect(s, x + 28, 875, 90, 30, fill=BG_GLASS_RGB, line=BORDER_RGB, radius=True)
    add_text(s, x + 28, 879, 90, 24, "PDF →",
             font=FONT_MONO, size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)

paper_card(158, "The Cost of Complexity", "In-progress draft",
           "Derives, from three sparse assumptions, the thermodynamic cost of any embodied representation — and a sharp critical ceiling where the partition function becomes the Riemann zeta function.")
paper_card(982, "Equilibrium from the Inside Out", "Exact Hamiltonian of Mean Force · In-progress draft",
           "An exact algebraic closure theorem: a criterion for when a subsystem of a larger physical system admits a compact local effective description — proved for a genuinely non-commuting case.")

add_counter(s, 3)
add_speaker_note(s, NOTES[2])

# ───── SLIDE 4 — One dynamics, many substrates ─────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG_PRIMARY)
add_kicker(s, 100, 60, "The unifying move")
add_text(s, 100, 110, 1720, 80, "One dynamics, many substrates.",
         font=FONT_DISPLAY, size=44, bold=True, color=TEXT_PRIMARY)
add_text(s, 100, 200, 1720, 80,
         [("Dynamical systems theory is the glue: biological, quantum, and neuromorphic learners are three physical realisations of the ", {"font": FONT_BODY, "size": 17, "color": TEXT_SECOND}),
          ("same operator algebra", {"font": FONT_BODY, "size": 17, "color": TEXT_PRIMARY, "bold": True}),
          (". Embodiment then imposes constraints — and constraints are where design principles come from.", {"font": FONT_BODY, "size": 17, "color": TEXT_SECOND})])

# Three substrate cards
def sub_card(x, label, title, foot, color):
    add_rect(s, x, 320, 555, 180, fill=BG_CARD_RGB, line=BORDER_RGB, line_w=1, radius=True)
    add_text(s, x + 24, 336, 500, 28, "— " + label.upper(),
             font=FONT_MONO, size=14, color=TEXT_MUTED, tracking=140)
    add_text(s, x + 24, 370, 500, 36, title,
             font=FONT_DISPLAY, size=22, bold=True, color=TEXT_PRIMARY)
    # mini "trace" — a thin rule + colored line
    add_rect(s, x + 24, 440, 500, 1, fill=BORDER_RGB)
    add_rect(s, x + 24, 438, 500, 3, fill=color)
    add_text(s, x + 24, 462, 500, 28, foot,
             font=FONT_BODY, size=14, italic=True, color=TEXT_SECOND)

sub_card(100,  "Biology",      "Spiking cortex",         "irregular spikes · sparse · cheap", GREEN)
sub_card(682,  "Quantum",      "Open-system evolution",  "coherent · non-commuting · lossy", BLUE)
sub_card(1264, "Neuromorphic", "Memristive reservoir",   "analog state · hysteretic · embodied", ACCENT)

# Glue bar
add_rect(s, 100, 580, 1720, 110, fill=RGBColor(0x1a,0x18,0x14), line=ACCENT, line_w=1, radius=True)
add_text(s, 130, 610, 280, 36, "— THE GLUE",
         font=FONT_MONO, size=14, color=ACCENT, bold=True, tracking=160)
add_text(s, 130, 638, 280, 40,
         [("∂", {"font": FONT_MONO, "size": 22, "color": TEXT_PRIMARY}),
          ("t", {"font": FONT_MONO, "size": 13, "color": TEXT_PRIMARY}),
          (" ρ = ", {"font": FONT_MONO, "size": 22, "color": TEXT_PRIMARY}),
          ("𝓛", {"font": FONT_MONO, "size": 22, "color": ACCENT, "bold": True}),
          ("[ρ]", {"font": FONT_MONO, "size": 22, "color": TEXT_PRIMARY})])
# vertical separator
add_rect(s, 470, 605, 1, 70, fill=ACCENT)
add_text(s, 500, 610, 1300, 70,
         "A single operator-algebraic representation — the same linear evolution over a lifted observable space, whatever the hardware underneath.",
         font=FONT_BODY, size=17, color=TEXT_SECOND)

# Two constraints
def constraint(x, tag, title, desc_runs):
    add_rect(s, x, 760, 850, 200, fill=BG_CARD_RGB, line=BORDER_RGB, line_w=1, radius=True)
    add_text(s, x + 28, 778, 800, 24, tag,
             font=FONT_MONO, size=13, color=TEXT_MUTED, tracking=140)
    add_text(s, x + 28, 808, 800, 40, title,
             font=FONT_DISPLAY, size=22, bold=True, color=TEXT_PRIMARY)
    add_text(s, x + 28, 854, 800, 100, desc_runs,
             font=FONT_BODY, size=15, color=TEXT_SECOND)

constraint(100, "CONSTRAINT · 01", "Thermodynamic cost",
           [("Every embodied representation pays a price in free energy. The ", {"font": FONT_BODY, "size": 15, "color": TEXT_SECOND}),
            ("Cost of Complexity", {"font": FONT_BODY, "size": 15, "color": ACCENT, "bold": True}),
            (" fixes that price from three sparse assumptions — and sets a critical ceiling.", {"font": FONT_BODY, "size": 15, "color": TEXT_SECOND})])
constraint(970, "CONSTRAINT · 02", "Algebraic closure",
           [("Not every substrate admits a compact local description of what it does. ", {"font": FONT_BODY, "size": 15, "color": TEXT_SECOND}),
            ("Equilibrium from the Inside Out", {"font": FONT_BODY, "size": 15, "color": ACCENT, "bold": True}),
            (" gives the exact criterion for when it does.", {"font": FONT_BODY, "size": 15, "color": TEXT_SECOND})])

add_footer(s, meta=""); add_counter(s, 4)
add_speaker_note(s, NOTES[3])

# ───── SLIDE 5 — Public writing ─────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG_PRIMARY)
add_kicker(s, 120, 130, "Explicability as deliverable")
add_text(s, 120, 180, 1700, 220,
         [("Public writing is\n", {"font": FONT_DISPLAY, "size": 60, "bold": True, "color": TEXT_PRIMARY}),
          ("part of the work.", {"font": FONT_DISPLAY, "size": 60, "bold": True, "color": TEXT_PRIMARY})])
add_text(s, 120, 410, 1500, 80,
         [("A synthesis which cannot be stated across the fields it unifies ", {"font": FONT_DISPLAY, "size": 24, "italic": True, "color": TEXT_SECOND}),
          ("hasn't unified them", {"font": FONT_DISPLAY, "size": 24, "color": ACCENT, "bold": True}),
          (".", {"font": FONT_DISPLAY, "size": 24, "italic": True, "color": TEXT_SECOND})])

def writing_card(x, label, title, blurb, meta_runs):
    add_rect(s, x, 580, 820, 380, fill=BG_CARD_RGB, line=BORDER_RGB, line_w=1, radius=True)
    add_text(s, x + 40, 610, 760, 28, label.upper(),
             font=FONT_MONO, size=13, bold=True, color=ACCENT, tracking=160)
    add_text(s, x + 40, 650, 760, 80, title,
             font=FONT_DISPLAY, size=32, bold=True, color=TEXT_PRIMARY)
    add_text(s, x + 40, 750, 760, 140, blurb,
             font=FONT_BODY, size=17, color=TEXT_SECOND)
    add_rect(s, x + 40, 905, 740, 1, fill=BORDER_RGB)
    add_text(s, x + 40, 915, 760, 32, meta_runs,
             font=FONT_MONO, size=14, color=TEXT_MUTED)

writing_card(120, "Ongoing essay practice", "Little Learning Machine",
             "Essays on physics, machine intelligence, and the occasionally unbearable nature of modern theory. Where Synthetics is drafted in public.",
             [("Substack", {"font": FONT_MONO, "size": 14, "color": ACCENT, "bold": True}),
              (" — littlelearningmachine.substack.com", {"font": FONT_MONO, "size": 14, "color": TEXT_MUTED})])
writing_card(980, "Essay prize", "Spanspermia",
             "\"Does Life Come from Outer Hilbert Space?\" — an argument that informational structure, not chemistry, is the substrate question worth asking.",
             [("2nd place · FQxI", {"font": FONT_MONO, "size": 14, "color": ACCENT, "bold": True}),
              (" Essay Competition", {"font": FONT_MONO, "size": 14, "color": TEXT_MUTED})])

add_footer(s, meta=""); add_counter(s, 5)
add_speaker_note(s, NOTES[4])

# ───── SLIDE 6 — About + URLs ─────
s = prs.slides.add_slide(blank); set_slide_bg(s, BG_PRIMARY)
add_kicker(s, 100, 100, "About")

# Portrait
if PORTRAIT.exists():
    add_rect(s, 100, 180, 520, 540, fill=BG_CARD_RGB, line=BORDER_RGB, line_w=1, radius=True)
    s.shapes.add_picture(str(PORTRAIT), E(120), E(200), width=E(480), height=E(500))
    add_text(s, 130, 715, 300, 24, "— THE APPLICANT",
             font=FONT_MONO, size=12, color=TEXT_MUTED, tracking=160)

# Right column
add_text(s, 700, 200, 1200, 110,
         [("Gerard ", {"font": FONT_DISPLAY, "size": 70, "bold": True, "color": TEXT_PRIMARY}),
          ("M", {"font": FONT_DISPLAY, "size": 70, "bold": True, "color": ACCENT}),
          ("c", {"font": FONT_DISPLAY, "size": 36, "bold": True, "color": ACCENT}),
          ("Caul.", {"font": FONT_DISPLAY, "size": 70, "bold": True, "color": ACCENT})])
add_text(s, 700, 320, 1200, 40,
         "Theoretical physicist · Loose canon",
         font=FONT_DISPLAY, size=22, italic=True, color=TEXT_SECOND)
add_text(s, 700, 380, 1200, 220,
         [("I work at the intersection of ", {"font": FONT_BODY, "size": 18, "color": TEXT_PRIMARY}),
          ("quantum theory, thermodynamics, and machine intelligence", {"font": FONT_BODY, "size": 18, "color": ACCENT, "bold": True}),
          (" — the three fields Synthetics proposes to unify. Twenty-five-plus papers across quantum control, open quantum systems, reservoir computing, superoscillations, and Koopman methods. I also write essays that are too long, and compose music that makes people wince.",
           {"font": FONT_BODY, "size": 18, "color": TEXT_PRIMARY})])

# Meta row
add_text(s, 700, 620, 1200, 32,
         [("PhD", {"font": FONT_MONO, "size": 14, "color": TEXT_PRIMARY, "bold": True}),
          (" · King's College London, 2019      ", {"font": FONT_MONO, "size": 14, "color": TEXT_MUTED}),
          ("FQxI", {"font": FONT_MONO, "size": 14, "color": TEXT_PRIMARY, "bold": True}),
          (" · 2nd place, 2024      ", {"font": FONT_MONO, "size": 14, "color": TEXT_MUTED}),
          ("Postdoctoral", {"font": FONT_MONO, "size": 14, "color": TEXT_PRIMARY, "bold": True}),
          (" · Tulane", {"font": FONT_MONO, "size": 14, "color": TEXT_MUTED})])

# URL row
add_rect(s, 700, 800, 1100, 1, fill=BORDER_RGB)
def url_cell(x, tag, base, accent):
    add_text(s, x, 820, 360, 24, tag,
             font=FONT_MONO, size=12, color=TEXT_MUTED, tracking=160)
    add_text(s, x, 850, 360, 36,
             [(base, {"font": FONT_DISPLAY, "size": 20, "color": TEXT_PRIMARY, "bold": True}),
              (accent, {"font": FONT_DISPLAY, "size": 20, "color": ACCENT, "bold": True})])

url_cell(700,  "01 · PHYSICS", "gmccaul.co.uk", "/physics.html")
url_cell(1080, "02 · ESSAYS",  "littlelearningmachine", ".substack.com")
url_cell(1460, "03 · WRITING", "gmccaul.co.uk", "/writing.html")
# vertical separators
add_rect(s, 1060, 820, 1, 80, fill=BORDER_RGB)
add_rect(s, 1440, 820, 1, 80, fill=BORDER_RGB)

add_footer(s, brand="Gerard McCaul", meta="Astera — Residency application · 2026")
add_counter(s, 6)
add_speaker_note(s, NOTES[5])

prs.save(OUT)
print(f"Wrote {OUT}")
