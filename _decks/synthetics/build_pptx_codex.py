"""Build a high-fidelity PowerPoint from rendered screenshots of the HTML deck.

This keeps the browser-rendered look of the slides by using the 1920x1080 PNG
captures as full-slide artwork, while still carrying over the speaker notes from
the HTML source into the PowerPoint notes pages.
"""

from __future__ import annotations

import argparse
import html
import json
import re
import subprocess
import sys
from pathlib import Path

from pptx import Presentation
from pptx.util import Emu

HERE = Path(__file__).resolve().parent
INDEX_HTML = HERE / "index.html"
SCREENSHOT_DIR = HERE / "slide-screenshots"
SCREENSHOT_SCRIPT = HERE / "screenshot_slides.mjs"
DEFAULT_OUT = HERE / "Synthetics Deck - Codex.pptx"

SLIDE_WIDTH = Emu(12192000)
SLIDE_HEIGHT = Emu(6858000)

SECTION_LABEL_RE = re.compile(r"<section\b[^>]*\bdata-label=\"([^\"]+)\"", re.IGNORECASE)
NOTES_RE = re.compile(
    r"<script\b[^>]*\bid=\"speaker-notes\"[^>]*>(.*?)</script>",
    re.IGNORECASE | re.DOTALL,
)


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description="Package the rendered Synthetics HTML deck into a PPTX."
    )
    parser.add_argument(
        "--html",
        type=Path,
        default=INDEX_HTML,
        help="Path to the deck HTML file.",
    )
    parser.add_argument(
        "--screenshots",
        type=Path,
        default=SCREENSHOT_DIR,
        help="Directory containing slide-XX.png screenshots.",
    )
    parser.add_argument(
        "--out",
        type=Path,
        default=DEFAULT_OUT,
        help="Output PPTX path.",
    )
    parser.add_argument(
        "--skip-render",
        action="store_true",
        help="Do not try to regenerate screenshots if any are missing.",
    )
    return parser.parse_args()


def load_html_metadata(index_html: Path) -> tuple[list[str], list[str]]:
    raw_html = index_html.read_text(encoding="utf-8")

    labels = [html.unescape(label) for label in SECTION_LABEL_RE.findall(raw_html)]

    notes_match = NOTES_RE.search(raw_html)
    if not notes_match:
        raise ValueError(f"Could not find #speaker-notes JSON in {index_html}")

    notes = json.loads(notes_match.group(1).strip())
    if not isinstance(notes, list):
        raise ValueError("Speaker notes JSON must be an array of note strings.")

    return labels, [str(note) for note in notes]


def screenshot_paths(screenshot_dir: Path) -> list[Path]:
    return sorted(
        path
        for path in screenshot_dir.glob("slide-*.png")
        if path.is_file()
    )


def ensure_screenshots(
    screenshot_dir: Path,
    expected_count: int,
    *,
    skip_render: bool,
) -> list[Path]:
    images = screenshot_paths(screenshot_dir)
    if len(images) == expected_count:
        return images

    if skip_render:
        raise FileNotFoundError(
            f"Expected {expected_count} screenshots in {screenshot_dir}, found {len(images)}."
        )

    if not SCREENSHOT_SCRIPT.exists():
        raise FileNotFoundError(f"Missing screenshot renderer: {SCREENSHOT_SCRIPT}")

    try:
        subprocess.run(
            ["node", str(SCREENSHOT_SCRIPT)],
            cwd=str(HERE),
            check=True,
        )
    except FileNotFoundError as exc:
        raise RuntimeError(
            "Node.js was not found, so screenshots could not be regenerated."
        ) from exc
    except subprocess.CalledProcessError as exc:
        raise RuntimeError(
            f"Screenshot rendering failed with exit code {exc.returncode}."
        ) from exc

    images = screenshot_paths(screenshot_dir)
    if len(images) != expected_count:
        raise RuntimeError(
            f"Expected {expected_count} screenshots after rendering, found {len(images)}."
        )
    return images


def build_presentation(images: list[Path], labels: list[str], notes: list[str]) -> Presentation:
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    blank = prs.slide_layouts[6]

    props = prs.core_properties
    props.title = "Synthetics Deck - Codex conversion"
    props.subject = "HTML slide deck conversion"
    props.author = "OpenAI Codex"
    props.comments = "Built from browser-rendered slide screenshots with imported speaker notes."

    for index, image_path in enumerate(images):
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(
            str(image_path),
            0,
            0,
            width=prs.slide_width,
            height=prs.slide_height,
        )

        label = labels[index] if index < len(labels) else f"Slide {index + 1}"
        note = notes[index] if index < len(notes) else ""
        note_parts = [label]
        if note.strip():
            note_parts.extend(["", note.strip()])
        slide.notes_slide.notes_text_frame.text = "\n".join(note_parts)

    return prs


def main() -> int:
    args = parse_args()

    labels, notes = load_html_metadata(args.html)
    expected_count = max(len(labels), len(notes))
    if expected_count == 0:
        raise ValueError("No slides were discovered in the HTML deck.")

    images = ensure_screenshots(
        args.screenshots,
        expected_count,
        skip_render=args.skip_render,
    )

    prs = build_presentation(images, labels, notes)
    args.out.parent.mkdir(parents=True, exist_ok=True)
    prs.save(args.out)

    print(f"Wrote {args.out}")
    print(f"Slides: {len(images)}")
    return 0


if __name__ == "__main__":
    try:
        raise SystemExit(main())
    except Exception as exc:  # pragma: no cover - CLI error path
        print(f"Error: {exc}", file=sys.stderr)
        raise SystemExit(1)
